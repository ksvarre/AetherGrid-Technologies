import os
from docx import Document
from pptx import Presentation
from openpyxl import Workbook

def ensure_dirs():
    os.makedirs("data/transcripts", exist_ok=True)
    os.makedirs("data/documents", exist_ok=True)
    os.makedirs("data/db", exist_ok=True)
    os.makedirs("scripts", exist_ok=True)

def generate_transcripts():
    transcripts = [
        {
            "filename": "transcript_2026_01_12_quantum_kickoff.md",
            "content": """---
date: 2026-01-12
attendees: Dr. Elena Rostova, Marcus Vance, Sarah Chen
facilitator: Sarah Chen
domain: Project Quantum
priority: Medium
---
# Meeting Transcript: Project Quantum Kickoff

**Sarah Chen**: Thanks for jumping in, everyone. We are kicking off Project Quantum today. The objective is to build our next-generation load forecasting model. Elena, where do we stand on the algorithm design?

**Dr. Elena Rostova**: Yes, hello. We are moving away from traditional auto-regressive models. I am planning a hybrid neural network architecture that integrates temporal fusion transformers with LSTM nodes. The target is to predict regional grid loads 24 hours in advance with a Mean Absolute Error (MAE) under 1.5 MW. 

**Marcus Vance**: That sounds solid, Elena. But remember, the model has to ingest smart-meter telemetry in real-time. Can we run this prediction loop on our standard cloud instances, or do we need dedicated GPU nodes?

**Dr. Elena Rostova**: The training phase definitely requires GPU acceleration, likely NVIDIA A10G instances. But the inference stage will be lightweight. We can run inference on standard CPU instances. We just need to ensure the database can stream substation load records without blocking.

**Sarah Chen**: Excellent. I have client meetings with Texas Electric next week, and they are extremely excited about the load prediction accuracy. 
- **Decision**: Elena will lead the core training pipeline using weather forecast arrays.
- **Action Item**: Marcus to provision a PostgreSQL staging database with timescale DB extensions by January 20th.
"""
        },
        {
            "filename": "transcript_2026_01_28_helium_sensor_issue.md",
            "content": """---
date: 2026-01-28
attendees: Marcus Vance, Dr. Elena Rostova, David Kross
facilitator: Marcus Vance
domain: Project Helium
priority: High
---
# Technical Alignment: Helium Sub-station Node Overheating

**Marcus Vance**: Okay team, we have a major blocker. The new Project Helium prototype nodes deployed at the Oak Creek sub-station are running extremely hot. During peak midday operations, the core enclosure temperatures reached 91°C, and the CPU throttled down to 800 MHz. This is causing sub-second latency spikes in the grid balancing loop.

**David Kross**: That's bad. The lightweight K3s Kubernetes pods running on those edge nodes are crash-looping because of the latency timeouts. Marcus, is this a physical fan failure or computational over-allocation?

**Marcus Vance**: It's a combination. The fan curve in our current firmware was locked to a low-RPM profile to keep noise levels down for residential sub-stations. But the thermal passive dissipation of the aluminum chassis is not enough when Elena's local telemetry pre-processing scripts are running at 100% CPU.

**Dr. Elena Rostova**: Wait, my pre-processing scripts are optimized! They are written in C++ and only perform standard rolling average windowing. If they are pegging the CPU, it means the incoming stream buffer has an infinite loop or duplicate telemetry readings.

**David Kross**: I checked the logs. The buffer was indeed getting flooded with duplicate smart-meter logs because of a broadcast storm in the substation router.
- **Decision**: Marcus will modify the firmware fan curve to trigger maximum cooling when core temps exceed 75°C.
- **Action Item**: David to isolate the smart-meter subnet to stop the packet duplication storm by February 2nd.
"""
        },
        {
            "filename": "transcript_2026_02_15_horizon_municipal_sales.md",
            "content": """---
date: 2026-02-15
attendees: Sarah Chen, Amira Patel, Marcus Vance
facilitator: Sarah Chen
domain: Project Horizon
priority: Medium
---
# Sales Alignment: Horizon Municipal Pitch for City of Austin

**Sarah Chen**: Welcome, everyone. We are finalizing our municipal pitch for the City of Austin's green energy initiative. They want to integrate 50 local solar microgrids into their central grid. Amira, how does Project Horizon solve their main integration concern?

**Amira Patel**: The City of Austin is terrified of battery depletion during sudden cloud cover. With Project Horizon, we introduce community battery discharge coordination. Our microgrid controllers negotiate in real-time. If Substation A detects a drop in solar voltage, it requests discharge credits from adjacent Substation B within 15 milliseconds.

**Marcus Vance**: Yes, but the physical battery hardware must support rapid discharge cycles. If they use cheap lithium iron phosphate batteries, we will degrade the cell life within two years. We must recommend premium LTO (Lithium Titanate) cells for the high-frequency balancing nodes.

**Sarah Chen**: That's a great commercial point, Marcus. I will structure a premium hardware tier that includes the LTO batteries. Amira, can you create a simplified grid flow diagram showing the battery discharge handshake?

**Amira Patel**: Yes, I can map the handshake logic. It will show how battery controllers communicate over the Horizon cellular telemetry channel.
- **Decision**: Recommending LTO cells as the default for the high-frequency municipal tier.
- **Action Item**: Amira to deliver the Austin microgrid grid flow diagram by February 20th.
"""
        },
        {
            "filename": "transcript_2026_03_02_database_scaling_crisis.md",
            "content": """---
date: 2026-03-02
attendees: David Kross, Dr. Elena Rostova, Marcus Vance
facilitator: David Kross
domain: DevOps / Database
priority: High
---
# Incident Review: Database Scaling Lock during Load Surge

**David Kross**: Let's review what happened yesterday. At 14:05 CST, during a regional heatwave, our main PostgreSQL metrics table locked up completely. The write latency went from 4ms to 48 seconds. Elena's load forecasting engine went blind because it couldn't retrieve recent Smart-Meter logs.

**Dr. Elena Rostova**: Yes, it was a disaster. The model had to fall back to historical averages, which led to a 12% over-allocation of power grid reserves. This cost the municipal client roughly $45,000 in energy market spot-pricing penalties.

**David Kross**: The culprit was database table lock contention. We had 120 edge nodes attempting to write 10-millisecond telemetry records into a single non-partitioned table simultaneously. The disk I/O on the primary AWS RDS instance hit 100% capacity and stayed there.

**Marcus Vance**: Why aren't we using TimescaleDB timeseries hyper-tables? They are designed specifically for high-throughput time-series ingestion.

**David Kross**: We planned to upgrade last month, but the migration script had a syntax error. I postponed it to focus on the edge container deployments. 
- **Decision**: We will perform emergency database table partitioning and implement a Redis cache layer for raw edge inputs to smooth out write spikes.
- **Action Item**: David to execute the PostgreSQL TimescaleDB hypertable migration by next Sunday.
"""
        },
        {
            "filename": "transcript_2026_03_19_gridpulse_pricing.md",
            "content": """---
date: 2026-03-19
attendees: Sarah Chen, Marcus Vance, Dr. Elena Rostova
facilitator: Sarah Chen
domain: Product Commercials
priority: Low
---
# Product Alignment: GridPulse Tiered Pricing Redesign

**Sarah Chen**: We need to update the commercial model for GridPulse. Our municipal clients are pushing back on the flat monthly fee. They want a value-based model tied to active smart-meters.

**Marcus Vance**: A meter-based licensing fee makes sense. But we have to account for edge node hardware costs. For Project Helium edge controllers installed at their physical substations, we incur a hard hardware manufacturing cost of $2,400 per node. We cannot bundle that into a pure software subscription.

**Dr. Elena Rostova**: Correct. Additionally, running the high-frequency neural forecasting models for large cities utilizes substantial cloud computing GPU power. A small utility with 10,000 meters shouldn't pay the same infrastructure overhead as a metro area with 500,000 meters.

**Sarah Chen**: I agree. Let's create three tiers:
1. **Utility Tier (Starter)**: Up to 50,000 meters. Local offline TF-IDF forecasting, 2 Helium edge nodes included.
2. **Municipal Tier (Standard)**: Up to 200,000 meters. Full Quantum neural forecasting, up to 10 Helium edge nodes, standard SLAs.
3. **Metropolitan Tier (Premium)**: Unlimited meters. Multi-region redundancy, high-frequency active balancing, 15ms battery discharge handshakes.

- **Decision**: Adopt a hybrid meter-based software model combined with a flat hardware setup fee.
- **Action Item**: Sarah to draft the formal GridPulse Pricing matrix Excel sheet by March 25th.
"""
        },
        {
            "filename": "transcript_2026_04_05_substation_safety_audit.md",
            "content": """---
date: 2026-04-05
attendees: Marcus Vance, Amira Patel, David Kross
facilitator: Marcus Vance
domain: Safety & Compliance
priority: Medium
---
# Compliance Review: Substation Physical Access Safety SOP

**Marcus Vance**: Safety audit is next week, everyone. We need to formalize the Standard Operating Procedure (SOP) for when technicians enter a physical high-voltage substation to upgrade our Project Helium edge nodes.

**Amira Patel**: This is highly regulated. Technicians must wear Category 4 Arc Flash protection gear before entering any enclosure. And we have a strict rule: never work on active high-voltage cabinets alone. A secondary compliance safety officer must be present outside the safety perimeter.

**David Kross**: From a systems perspective, we need a digital lock out protocol. Before a technician opens a Helium edge enclosure door, they must trigger "Maintenance Mode" in the central control app. This safely stops active container operations, flushes telemetry queues, and disables remote firmware pushes.

**Marcus Vance**: Good point, David. That prevents a random remote firmware update from bricking the node mid-service.
- **Decision**: Incorporating the Maintenance Mode trigger into the official technician mobile application.
- **Action Item**: Amira to update the physical safety checklist PDF and submit it to the regulator by April 12th.
"""
        },
        {
            "filename": "transcript_2026_04_20_quantum_training_gaps.md",
            "content": """---
date: 2026-04-20
attendees: Dr. Elena Rostova, Sarah Chen, David Kross
facilitator: Dr. Elena Rostova
domain: Project Quantum
priority: High
---
# Algorithm Review: Weather Sensor Gaps in ML Training Data

**Dr. Elena Rostova**: I am seeing a worrying drop in our Project Quantum forecasting accuracy. The MAE for the North region spiked to 2.4 MW yesterday. The model is misinterpreting sudden wind speed changes, which led to incorrect wind turbine output predictions.

**Sarah Chen**: Elena, is this an algorithm bug or a data ingestion issue?

**Dr. Elena Rostova**: It is a critical training data gap. The weather telemetry feed from our primary vendor went offline for 18 hours last week. Our database filled the missing wind speed column with zero values instead of interpolating the missing data. The model assumed there was a dead calm wind period, when in fact there was a 35-knot wind surge!

**David Kross**: That's a data engineering pipeline failure. Our ETL script should never insert raw zeros for missing values. It should calculate a rolling average or throw a warning flag so we don't feed corrupt vectors into the model.

**Dr. Elena Rostova**: Exactly. I need clean weather inputs.
- **Decision**: Update the ETL pipeline to use linear interpolation for data gaps under 2 hours, and flag gaps over 2 hours as system warnings.
- **Action Item**: David to rewrite the Python data ingestion script in our Kubernetes pipeline by Friday.
"""
        },
        {
            "filename": "transcript_2026_05_01_helium_antenna_planning.md",
            "content": """---
date: 2026-05-01
attendees: Marcus Vance, Amira Patel, Dr. Elena Rostova
facilitator: Marcus Vance
domain: Project Helium
priority: Medium
---
# Engineering Review: Substation Wireless Telemetry Antennas

**Marcus Vance**: We are designing the antenna bracket layouts for Project Helium edge nodes. We need reliable cellular backup connections because some municipal substations lack landline fiber links.

**Amira Patel**: Substations are a nightmare for electromagnetic interference (EMI). The high-voltage transformers radiate a lot of high-frequency noise. If we place the cellular antenna too close to the main transformer bank, the signal-to-noise ratio will collapse.

**Dr. Elena Rostova**: What spectrum are we using? If we use standard 2.4 GHz bands, the interference will be severe. We should stick to LTE Band 48 (CBRS) or dedicated 900 MHz industrial radio links.

**Marcus Vance**: We are standardizing on dual-band industrial modems that support CBRS and LTE Band 13 for long-distance rural coverage. We will mount the antenna on a 3-meter fiberglass extension pole to clear the metallic cages.
- **Decision**: Cellular antennas must be separated from primary transformer coils by a minimum distance of 5 meters.
- **Action Item**: Marcus to finalize the physical CAD brackets for the antenna mounts by May 10th.
"""
        },
        {
            "filename": "transcript_2026_05_10_horizon_microgrid_rules.md",
            "content": """---
date: 2026-05-10
attendees: Amira Patel, Sarah Chen, Marcus Vance
facilitator: Amira Patel
domain: Project Horizon
priority: Medium
---
# Policy Design: Community Battery Discharge Regulations

**Amira Patel**: We need to set the software rules for Project Horizon community battery discharging. If a local residential battery system is enrolled, how much capacity can the public utility draw during a grid alert?

**Sarah Chen**: Our customer agreement says we can draw up to 40% of their battery capacity during an official ERCOT grid alert, but we must guarantee their battery is never depleted below a 20% state of charge (SoC) buffer.

**Marcus Vance**: Yes, and we must check battery health before drawing. If the battery cell temperature is above 45°C or showing a high internal resistance spike, the microgrid controller must bypass that node to avoid damaging client hardware.

**Amira Patel**: Agreed. I will write these boundary constraints into the Horizon firmware rules. The controller will query local temperature and SoC registers before accepting a discharge command.
- **Decision**: Minimum customer SoC buffer locked at 20%. Bypass node if temperature exceeds 45°C.
- **Action Item**: Amira to implement these constraints in the battery discharge state machine by May 15th.
"""
        },
        {
            "filename": "transcript_2026_05_15_quarterly_roadmap_review.md",
            "content": """---
date: 2026-05-15
attendees: Sarah Chen, Dr. Elena Rostova, Marcus Vance, David Kross
facilitator: Sarah Chen
domain: Product Roadmap
priority: Medium
---
# Strategy Alignment: Q2 Product and Tech Roadmap Review

**Sarah Chen**: It's time for our quarterly roadmap review. We have major milestones coming up. Elena, what is the status of the Project Quantum model integration?

**Dr. Elena Rostova**: The model core is complete. We finished training v2.0 on the weather histories and smart-meter dataset. We achieved a 1.15 MW MAE in local validation, which beats our 1.5 MW target. We are ready to deploy to the production Kubernetes environment next week.

**David Kross**: That fits my schedule. I have finalized the edge orchestration spec for Project Helium. We will run the lightweight container packages on Marcus's v2 chassis nodes using K3s clusters.

**Marcus Vance**: Perfect. The physical hardware nodes are assembled and have passed the safety audit. We are shipping the first batch of 15 nodes to utility stations on June 1st.
- **Decision**: Q2 milestones are locked. Quantum model goes live on May 24th.
- **Action Item**: David to configure the CI/CD pipeline for automated model rollouts by May 20th.
"""
        },
        {
            "filename": "transcript_2026_05_18_devops_kubernetes_edge.md",
            "content": """---
date: 2026-05-18
attendees: David Kross, Marcus Vance, Dr. Elena Rostova
facilitator: David Kross
domain: DevOps / Helium
priority: Medium
---
# DevOps Alignment: Edge Cluster Container Orchestration

**David Kross**: We are aligning on container management for the Project Helium edge nodes. Since these nodes are physically distributed, standard heavy Kubernetes is out of the question. We are using Rancher K3s.

**Marcus Vance**: David, what happens if an edge node loses its cell connection? Does it stop balancing the local substation power grid?

**David Kross**: No. The local grid balancing loop runs as a standalone C++ micro-service that talks directly to the micro-controller over Modbus. It operates independently of the K3s control plane. If connection is lost, it continues autonomous operations and buffers logs locally. Once the cellular link recovers, the container syncs the metrics database.

**Dr. Elena Rostova**: That is excellent. I am packaging the local predictive model in a lightweight WASM container. This will allow it to run on the node with less than 256MB of RAM.
- **Decision**: Standardizing on Rancher K3s with local autonomous offline loop execution.
- **Action Item**: David to publish the edge deployment Docker templates by May 22nd.
"""
        },
        {
            "filename": "transcript_2026_05_19_emergency_load_simulation.md",
            "content": """---
date: 2026-05-19
attendees: Dr. Elena Rostova, Amira Patel, David Kross
facilitator: Dr. Elena Rostova
domain: AI Testing / Grid Simulation
priority: High
---
# Simulation Analysis: Solar Flare EMP Grid Resilience

**Dr. Elena Rostova**: Yesterday we simulated a solar flare EMP grid emergency. We injected synthetic voltage fluctuations into the Project Horizon simulation model to see if our AI controllers could prevent cascading power blackouts.

**Amira Patel**: The results were fascinating. The microgrid controllers detected the anomaly within 8 milliseconds. They successfully coordinated the community battery discharge loop, injecting 14 MW of buffer power to stabilize local municipal voltage lines.

**David Kross**: Yes, but the database logs showed a massive queue overflow. The telemetry workers on the database container couldn't keep up with the millisecond write rate. We lost about 3% of the metrics logs during the surge.

**Dr. Elena Rostova**: Losing metrics logs is acceptable during an EMP emergency, but we must protect grid stability. The neural forecaster maintained an active load-balancing state despite the data loss.
- **Decision**: Solar flare grid balancing simulation has officially passed.
- **Action Item**: David to optimize the database telemetry write buffer pool by May 23rd.
"""
        }
    ]

    import re
    for t in transcripts:
        filepath = f"data/transcripts/{t['filename']}"
        # Programmatically strip frontmatter block to ensure transcripts on disk are PURE dialogue
        pure_content = re.sub(r"^---\r?\n([\s\S]*?)\r?\n---\r?\n", "", t["content"])
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(pure_content)
    print(f"Generated {len(transcripts)} meeting transcripts.")

def generate_docx():
    docs = [
        {
            "filename": "quantum_ml_forecasting_spec.docx",
            "title": "Technical Specification: Project Quantum Neural Load Forecasting Model",
            "author": "Dr. Elena Rostova",
            "date": "2026-04-12",
            "content": [
                ("1. Executive Summary", "Project Quantum represents AetherGrid's core artificial intelligence model designed to predict grid loads. By utilizing neural networks, we predict regional electrical usage with high accuracy, enabling utility operators to balance grids and avoid expensive blackouts."),
                ("2. Model Architecture", "The forecasting engine uses a hybrid machine learning model integrating Temporal Fusion Transformers (TFT) with deep LSTM recurrent nodes. Weather history arrays (ambient temperature, wind speed, solar irradiance) are combined with active smart-meter usage statistics to generate a rolling 24-hour load prediction profile."),
                ("3. Accuracy Metrics & Performance", "The target specification for the predictive model is a Mean Absolute Error (MAE) threshold under 1.5 MW. In physical validation testing on North-region municipal datasets, the model achieved a stable MAE of 1.15 MW. High accuracy forecasting is critical during extreme weather events where supply-demand spikes occur."),
                ("4. Fail-Safe Operations", "In the event of a weather telemetry feed failure, the system throws an automated warning flag. The model falls back to a rolling 7-day average load curve, maintaining an active but low-confidence load-balancing profile until telemetry recovers.")
            ]
        },
        {
            "filename": "helium_edge_node_hardware_v2.docx",
            "title": "Hardware Design Specification: Project Helium Edge Sub-Station Controller v2",
            "author": "Marcus Vance",
            "date": "2026-03-15",
            "content": [
                ("1. Scope", "This document details the chassis design, electrical safety envelopes, and thermal performance parameters of the Project Helium edge compute node v2, designed for deployment at municipal electrical sub-stations."),
                ("2. Enclosure and Chassis", "The Helium v2 chassis is fabricated from extruded anodized aluminum (Grade 6061-T6) featuring external passive heat dissipation fins. The enclosure matches NEMA 4X weatherproofing standards to prevent dust ingress and moisture condensation."),
                ("3. Thermal Limits and Active Cooling", "Under standard operating conditions, the passive thermal fins maintain a safe core enclosure temperature. However, under peak computational load (100% CPU allocation), internal temperatures can rise to 91°C. To prevent thermal throttling, we have integrated a high-performance brushless active cooling fan. The firmware is configured to engage the fan at 100% speed when the core enclosure temperature exceeds 75°C. Safe operations require keeping the temperature below 85°C."),
                ("4. Cellular Telemetry brackets", "Helium edge nodes are equipped with industrial LTE modems supporting Band 48 (CBRS) and Band 13 (700 MHz). Telemetry antennas must be separated from primary transformer coils by a minimum distance of 5 meters on a fiberglass extension pole to prevent electromagnetic signal noise.")
            ]
        },
        {
            "filename": "devops_edge_infrastructure_scaling.docx",
            "title": "Systems Specification: Kubernetes Container Scaling at the Edge",
            "author": "David Kross",
            "date": "2026-05-02",
            "content": [
                ("1. Infrastructure Design", "Deploying high-frequency telemetry parsers and neural forecasting models at distributed sub-stations requires a robust, lightweight container orchestration plane. We have selected Rancher K3s as our standard lightweight Kubernetes distribution."),
                ("2. Node Orchestration and Memory Budgets", "Each physical Helium edge host node allocates a maximum budget of 4 CPUs and 8 GB of RAM. The local telemetry parser runs as a core daemonset, and the neural inference engine is packaged into a lightweight WebAssembly (WASM) container. This limits the total memory footprint of the active predictive nodes to less than 256MB of RAM."),
                ("3. Database Synchronization and Offline Buffering", "In case of a cellular network dropout, the local grid balancing containers continue executing the grid feedback loop autonomously over local Modbus channels. Telemetry logs are buffered locally in an on-disk SQLite queue. Once the connection is re-established, the edge node flushes the buffer to the central PostgreSQL hyper-table.")
            ]
        },
        {
            "filename": "substation_operations_compliance_sop.docx",
            "title": "Compliance Standard Operating Procedure: Sub-station Enclosure Maintenance",
            "author": "Amira Patel",
            "date": "2026-03-10",
            "content": [
                ("1. Compliance Mandate", "Technicians entering physical high-voltage substations to service Project Helium edge nodes must follow strict safety mandates. High-voltage utility equipment poses severe arc flash and electrocution hazards."),
                ("2. Protective Equipment & Safety Teams", "Before entering any enclosure, technicians must equip themselves with Category 4 Arc Flash protection gear (face shield, fire-resistant clothing, and insulated gloves). To ensure physical safety, never work on active high-voltage cabinets alone. A secondary compliance safety officer must be present outside the safety perimeter at all times."),
                ("3. Maintenance Mode Digital Protocol", "Before opening a physical edge node enclosure, the technician must trigger 'Maintenance Mode' in the central technician control app. This disables remote firmware updates, safely pauses active K3s container operations, and flushes database write queues to prevent data corruption.")
            ]
        }
    ]

    for d in docs:
        doc = Document()
        doc.add_heading(d["title"], level=1)
        doc.add_paragraph(f"Author: {d['author']}")
        doc.add_paragraph(f"Date: {d['date']}")
        doc.add_paragraph("---")
        for heading, body in d["content"]:
            doc.add_heading(heading, level=2)
            doc.add_paragraph(body)
        filepath = f"data/documents/{d['filename']}"
        doc.save(filepath)
    print(f"Generated {len(docs)} Word documents.")

def generate_pptx():
    decks = [
        {
            "filename": "gridpulse_product_roadmap_q2_2026.pptx",
            "title": "GridPulse Product Roadmap: Q2 2026",
            "author": "Sarah Chen",
            "date": "2026-05-01",
            "slides": [
                ("GridPulse Q2 Roadmap", "AetherGrid Technologies core software roadmap. Author: Sarah Chen. Date: May 1st, 2026."),
                ("Feature Releases", "- Project Quantum Integration: High-accuracy AI neural load forecasting goes live in the production SaaS console.\n- Multi-Grid Solar Battery balancing: Real-time community solar discharge handshake support.\n- Remote Substation Maintenance: Integrated maintenance mode inside the grid operator app."),
                ("Licensing and Commercial Tiers", "- Utility Tier: For small operators up to 50,000 smart-meters. Uses local forecasting, includes 2 edge nodes.\n- Municipal Tier: For mid-size operators up to 200,000 smart-meters. Full neural forecasting, up to 10 edge nodes.\n- Metropolitan Tier: Unlimited meters, multi-region failover, 15ms battery balancing SLA.")
            ]
        },
        {
            "filename": "horizon_decentralized_microgrid_pitch.pptx",
            "title": "Project Horizon: Decentralized Microgrids",
            "author": "Amira Patel",
            "date": "2026-02-10",
            "slides": [
                ("Project Horizon Pitch Deck", "Municipal Smart Grid Integration. Author: Amira Patel. Date: Feb 10th, 2026."),
                ("Real-time Power Negotiator", "- Real-time battery discharge negotiation allows local substations to exchange power spikes within 15 milliseconds.\n- Maximizes community solar usage and protects municipal grids from sudden load surges or cloud cover brownouts."),
                ("Hardware Battery Recommendations", "- Lithium Titanate (LTO) battery cells are the default recommendation for municipal tier customers. LTO supports rapid, high-frequency charge-discharge cycles without degrading cell capacity, providing a 10-year service life compared to 2 years for standard LFP cells.")
            ]
        },
        {
            "filename": "helium_substation_deployment_plan.pptx",
            "title": "Project Helium: Edge Node Substation Deployment Plan",
            "author": "Marcus Vance",
            "date": "2026-04-18",
            "slides": [
                ("Helium Substation Deployment", "Edge hardware rollout plan. Author: Marcus Vance. Date: April 18th, 2026."),
                ("Deployment Stages", "- Phase 1 (June 1st): Ship initial batch of 15 Helium edge nodes to local smart substations.\n- Phase 2 (August 15th): Install rural antennas and cellular CBRS brackets on fiberglass poles.\n- Phase 3 (November 1st): Full remote Kubernetes integration via Rancher K3s control plane."),
                ("Substation Physical Separations", "- Safety compliance rule: Antenna mounts must maintain a minimum physical distance of 5 meters from primary high-voltage transformers to ensure signal-to-noise ratio is protected from EMF noise.")
            ]
        },
        {
            "filename": "quantum_load_prediction_demo.pptx",
            "title": "Project Quantum: Predictive AI Model Demo",
            "author": "Dr. Elena Rostova",
            "date": "2026-05-12",
            "slides": [
                ("Project Quantum AI Forecaster", "Temporal Fusion Transformer Load Forecasting. Author: Dr. Elena Rostova. Date: May 12th, 2026."),
                ("Model Performance Demo", "- Validation testing on regional smart-meter datasets achieved a stable MAE of 1.15 MW.\n- High-accuracy load predictions enable local utility managers to dynamically buy grid reserves ahead of time, saving up to 15% in daily power sourcing costs."),
                ("Data Quality Safeguards", "- Core model requires clean weather and wind speed sensor feeds.\n- If weather feeds drop, the system flags a data warning and interpolates missing fields to prevent feeding zero vectors to the neural net.")
            ]
        }
    ]

    for d in decks:
        prs = Presentation()
        # Add Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = d["title"]
        subtitle.text = f"Author: {d['author']} | Date: {d['date']}"

        # Add Content Slides
        for s_title, s_text in d["slides"][1:]:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = s_title
            slide.placeholders[1].text = s_text

        filepath = f"data/documents/{d['filename']}"
        prs.save(filepath)
    print(f"Generated {len(decks)} PowerPoint presentations.")

def generate_xlsx():
    # 1. quantum_model_benchmarks_v1.xlsx
    wb1 = Workbook()
    ws1 = wb1.active
    ws1.title = "Model Accuracy Metrics"
    ws1.append(["Metadata Block", "", "", "", ""])
    ws1.append(["File: quantum_model_benchmarks_v1.xlsx", "", "", "", ""])
    ws1.append(["Author: Dr. Elena Rostova", "", "", "", ""])
    ws1.append(["Date: 2026-05-14", "", "", "", ""])
    ws1.append(["", "", "", "", ""])
    ws1.append(["Substation Region", "Model Type", "MAE Target (MW)", "MAE Tested (MW)", "Validation Status"])
    ws1.append(["Austin North", "Temporal Fusion Transformer", 1.5, 1.15, "Pass"])
    ws1.append(["Austin South", "Temporal Fusion Transformer", 1.5, 1.28, "Pass"])
    ws1.append(["West Texas Edge", "Lightweight LSTM", 1.5, 1.42, "Pass"])
    ws1.append(["Houston Metro", "Hybrid Transformer-LSTM", 1.5, 1.08, "Pass"])
    wb1.save("data/documents/quantum_model_benchmarks_v1.xlsx")

    # 2. helium_hardware_thermal_tests.xlsx
    wb2 = Workbook()
    ws2 = wb2.active
    ws2.title = "Helium Thermal Profile"
    ws2.append(["Metadata Block", "", "", "", "", ""])
    ws2.append(["File: helium_hardware_thermal_tests.xlsx", "", "", "", "", ""])
    ws2.append(["Author: Marcus Vance", "", "", "", "", ""])
    ws2.append(["Date: 2026-04-12", "", "", "", "", ""])
    ws2.append(["", "", "", "", "", ""])
    ws2.append(["Test ID", "Computational Load (%)", "Fan Speed Profile", "Ambient Temp (°C)", "Core Temp (°C)", "Safety Status"])
    ws2.append(["Test_01", 10, "Idle Passive", 25, 42, "Pass"])
    ws2.append(["Test_02", 50, "Quiet (40% RPM)", 25, 58, "Pass"])
    ws2.append(["Test_03", 100, "Maximum (100% RPM)", 25, 82, "Pass"])
    ws2.append(["Test_04", 100, "Throttled (Passive)", 25, 91, "Throttling Danger"])
    ws2.append(["Test_05", 100, "Maximum (100% RPM)", 40, 84, "Pass"])
    wb2.save("data/documents/helium_hardware_thermal_tests.xlsx")

    # 3. infrastructure_node_capacity_allocations.xlsx
    wb3 = Workbook()
    ws3 = wb3.active
    ws3.title = "K3s Pod Allocation"
    ws3.append(["Metadata Block", "", "", "", "", ""])
    ws3.append(["File: infrastructure_node_capacity_allocations.xlsx", "", "", "", "", ""])
    ws3.append(["Author: David Kross", "", "", "", "", ""])
    ws3.append(["Date: 2026-05-18", "", "", "", "", ""])
    ws3.append(["", "", "", "", "", ""])
    ws3.append(["Substation Node", "Edge Hostname", "CPUs Allocated", "Memory Allocated (GB)", "Active Containers", "CPU Usage (%)"])
    ws3.append(["Austin Sub-08", "austin-edge-08", 4, 8, 5, 42])
    ws3.append(["Austin Sub-12", "austin-edge-12", 4, 8, 6, 88])
    ws3.append(["Dallas Sub-01", "dallas-edge-01", 4, 8, 4, 18])
    ws3.append(["Houston Sub-24", "houston-edge-24", 4, 8, 5, 54])
    wb3.save("data/documents/infrastructure_node_capacity_allocations.xlsx")

    # 4. gridpulse_licensing_pricing_matrix.xlsx
    wb4 = Workbook()
    ws4 = wb4.active
    ws4.title = "Pricing Matrix"
    ws4.append(["Metadata Block", "", "", "", "", ""])
    ws4.append(["File: gridpulse_licensing_pricing_matrix.xlsx", "", "", "", "", ""])
    ws4.append(["Author: Sarah Chen", "", "", "", "", ""])
    ws4.append(["Date: 2026-03-24", "", "", "", "", ""])
    ws4.append(["", "", "", "", "", ""])
    ws4.append(["Client Tier", "Utility Smart-Meters", "Base Monthly License ($)", "Per-Meter Charge ($)", "Included Helium Edge Nodes", "Muncipal SLA Level"])
    ws4.append(["Utility Starter", 50000, 1500, 0.05, 2, "Standard (Offline Forecasting)"])
    ws4.append(["Municipal Standard", 200000, 5000, 0.03, 10, "High-Priority Neural Forecasting"])
    ws4.append(["Metropolitan Premium", 500000, 12000, 0.02, 15, "15ms Real-time Battery Discharge"])
    wb4.save("data/documents/gridpulse_licensing_pricing_matrix.xlsx")

    print("Generated 4 Excel spreadsheets.")

if __name__ == "__main__":
    ensure_dirs()
    generate_transcripts()
    generate_docx()
    generate_pptx()
    generate_xlsx()
    print("Synthetic database completely constructed!")
